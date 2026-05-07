"""Generate docs/Presentation.pptx for the BirdCLEF 2026 milestone presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "docs" / "Presentation.pptx"

# Brand palette
NAVY = RGBColor(0x0B, 0x2D, 0x5B)
TEAL = RGBColor(0x1F, 0x6F, 0x6F)
ORANGE = RGBColor(0xE3, 0x6E, 0x21)
GREY = RGBColor(0x55, 0x5B, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF6, 0xF7, 0xF9)


def _set_run(run, text, *, size=18, bold=False, color=NAVY, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _title_only_layout(prs):
    return prs.slide_layouts[5]  # title only


def _blank_layout(prs):
    return prs.slide_layouts[6]


def _add_title(slide, text, *, size=32, color=NAVY):
    if slide.shapes.title is None:
        return
    title = slide.shapes.title
    title.text_frame.text = text
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(size)
    title.text_frame.paragraphs[0].runs[0].font.color.rgb = color
    title.text_frame.paragraphs[0].runs[0].font.bold = True


def _add_textbox(slide, left, top, width, height, *, fill=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    return tf


def _add_para(tf, text, *, size=18, bold=False, color=NAVY, bullet=True, indent=0):
    if tf.paragraphs[0].text == "" and bullet is False and not bold:
        # reuse first paragraph if empty
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.level = indent
    run = p.add_run()
    if bullet and not text.startswith(("•", "  ")):
        text = "• " + text
    _set_run(run, text, size=size, bold=bold, color=color)
    return p


def _new_slide(prs, layout_idx=5):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def _add_bullets(slide, bullets, *, left=0.5, top=1.6, width=12.5, height=5.5,
                 size=20, color=NAVY):
    tf = _add_textbox(slide, left, top, width, height)
    first = True
    for b in bullets:
        if isinstance(b, tuple):
            text, lvl = b
        else:
            text, lvl = b, 0
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = lvl
        run = p.add_run()
        prefix = "• " if lvl == 0 else "— "
        _set_run(run, prefix + text, size=size - lvl * 2, color=color)


def _add_table(slide, headers, rows, *, left=0.5, top=1.6, width=12.5, height=5.0,
               header_fill=NAVY, body_fill=LIGHT_BG):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        _set_run(run, h, size=16, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            color = NAVY
            text = str(val)
            if text.startswith("**") and text.endswith("**"):
                text = text.strip("*")
                run = p.add_run() if cell.text_frame.paragraphs[0].runs else p.add_run()
                _set_run(run, text, size=14, bold=True, color=ORANGE)
            else:
                _set_run(run, text, size=14, color=color)
            cell.fill.solid()
            cell.fill.fore_color.rgb = body_fill
    return table


def _add_footer(slide, text="BirdCLEF+ 2026 — Pantanal Acoustic Species ID"):
    tf = _add_textbox(slide, 0.4, 7.05, 13.0, 0.4)
    p = tf.paragraphs[0]
    run = p.add_run()
    _set_run(run, text, size=11, color=GREY, italic=True)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # widescreen
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title ----
    s = _new_slide(prs, layout_idx=6)  # blank
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tf = _add_textbox(s, 0.6, 2.6, 12, 1.2)
    p = tf.paragraphs[0]
    run = p.add_run()
    _set_run(run, "Acoustic Species Identification", size=44, bold=True, color=WHITE)

    tf = _add_textbox(s, 0.6, 3.6, 12, 0.8)
    p = tf.paragraphs[0]
    run = p.add_run()
    _set_run(run, "in the Brazilian Pantanal — BirdCLEF+ 2026", size=28, color=ORANGE)

    tf = _add_textbox(s, 0.6, 4.7, 12, 0.6)
    p = tf.paragraphs[0]
    run = p.add_run()
    _set_run(
        run,
        "A Heterogeneous Rank-Average Ensemble — Final Public LB 0.932",
        size=20,
        color=WHITE,
    )

    tf = _add_textbox(s, 0.6, 6.0, 12, 0.6)
    p = tf.paragraphs[0]
    run = p.add_run()
    _set_run(
        run,
        "Savitha Vijayarangan  ·  github.com/RamananVr/KaggleBirdClef2026",
        size=14,
        color=WHITE,
        italic=True,
    )

    # ---- Slide 2: Introduction ----
    s = _new_slide(prs)
    _add_title(s, "1. Introduction")
    _add_bullets(
        s,
        [
            "Goal: identify which of 234 species (birds, amphibians, mammals, reptiles, insects) is calling in 5-second windows of 1-minute Pantanal soundscapes",
            "Real-world impact: passive acoustic monitoring at scale enables biodiversity tracking in remote, inaccessible wetlands threatened by climate, agriculture, fire",
            ("Constraints: code competition — CPU only, no internet at submission time, ≤ 90 min runtime, hidden test data only mounted during scoring", 0),
            "Evaluation: macro-averaged ROC-AUC over classes with at least one positive in hidden test",
            "Daily quota: 5 submissions per UTC day; up to 2 picks for final private-leaderboard judgment",
            "Our best public LB: 0.932 — top of leaderboard 0.956",
        ],
        size=18,
    )
    _add_footer(s)

    # ---- Slide 3: Related Work ----
    s = _new_slide(prs)
    _add_title(s, "2. Related Work")
    _add_bullets(
        s,
        [
            "Foundation models for bioacoustics — Google Perch v2 (Cornell + Google) trained on a global bird-vocalization corpus; produces logits + 1536-d embeddings",
            ("→ Used as our strongest single branch", 1),
            "Distilled Sound Event Detection ensembles — multi-fold mel-spec models compressed to ONNX for CPU code-comp deployment",
            ("→ Tucker Arrants public 5-fold ensemble (CC-licensed)", 1),
            "Sequence modeling over Perch embeddings — selective state-space + cross-attention + KD (the public 0.943 architecture)",
            ("→ We attempted a stripped-down bi-LSTM variant — see Section 6", 1),
            "Rank-average ensembling — robust to scale mismatch across heterogeneous models",
            "Post-processing — file-level rank-aware scaling and delta-shift smoothing (cited as 2025 Rank-1/3 techniques)",
        ],
        size=17,
    )
    _add_footer(s)

    # ---- Slide 4: Data Sources ----
    s = _new_slide(prs)
    _add_title(s, "3.1 Data Sources")
    _add_table(
        s,
        ["Source", "Files", "Role"],
        [
            ("train_audio/ (Xeno-Canto + iNaturalist)", "35,549 OGGs", "Foreground species clips"),
            ("train_soundscapes/", "10,658 OGGs", "1-min Pantanal recordings"),
            ("train_soundscapes_labels.csv", "1,478 segments / 66 files", "Expert annotations (5-s segments)"),
            ("test_soundscapes/", "~600 OGGs (hidden)", "Mounted only at scoring re-run"),
            ("taxonomy.csv", "234 classes", "Maps primary_label → scientific name + class_name"),
            ("sample_submission.csv", "235 columns", "row_id + 234 species probabilities"),
        ],
        top=1.6,
        height=4.8,
    )
    _add_footer(s)

    # ---- Slide 5: Cleansing & Validation ----
    s = _new_slide(prs)
    _add_title(s, "3.2 Data Cleansing, Validation & Class Coverage")
    _add_bullets(
        s,
        [
            "Audio loaded with `soundfile` (libsndfile-backed) — avoids the recent torchaudio→torchcodec dependency",
            "Mean-pool stereo to mono; resample to 32 kHz (already in source); pad/truncate to 60 s",
            "NaN-safe mel-spectrogram with ε before dB conversion to keep silent inputs finite",
            "Critical join: 234 target classes × Perch's vocabulary",
            ("203 / 234 classes (87%) have a direct match by scientific_name", 1),
            ("3 classes get a genus-level proxy (max-pool over same-genus Perch entries)", 1),
            ("**28 classes (12%) have no Perch match — insect sonotypes like 47158son16**", 1),
            "→ This single observation drove our multi-branch design",
            "Validation: 5-fold stratified CV on XC+iNat; soundscape labels held out per-fold; LB-only feedback for post-processing knobs",
        ],
        size=17,
    )
    _add_footer(s)

    # ---- Slide 6: Transformation & Visualization ----
    s = _new_slide(prs)
    _add_title(s, "3.3 Data Transformation & Visualization")

    tf = _add_textbox(s, 0.5, 1.5, 6.0, 5.5)
    _add_para(tf, "Per-branch preprocessing:", size=18, bold=True, bullet=False)
    _add_para(tf, "Perch:  raw waveform [N, 160_000] @ 32 kHz", size=15, bullet=True)
    _add_para(tf, "→ ONNX inference returns 'label' + 'embedding'", size=14, bullet=False, indent=1)
    _add_para(tf, "SED:  log-mel spec (256 mels, n_fft=2048, hop=512)", size=15, bullet=True)
    _add_para(tf, "→ frame-level + clip-level outputs, sigmoid-averaged", size=14, bullet=False, indent=1)
    _add_para(tf, "EffNet:  log-mel spec (128 mels, hop=512)", size=15, bullet=True)
    _add_para(tf, "→ per-sample standardize, SpecAugment in training", size=14, bullet=False, indent=1)
    _add_para(tf, "Each branch produces a (12, 234) matrix per file", size=15, bullet=True, color=ORANGE)

    tf = _add_textbox(s, 7.0, 1.5, 6.0, 5.5, fill=LIGHT_BG)
    _add_para(tf, "Class distribution (234 total):", size=18, bold=True, bullet=False, color=NAVY)
    _add_para(tf, "Aves (birds)  ≈ 178", size=15, bullet=True)
    _add_para(tf, "Insecta  ≈ 30 (sonotype IDs)", size=15, bullet=True)
    _add_para(tf, "Amphibia  ≈ 20", size=15, bullet=True)
    _add_para(tf, "Reptilia + Mammalia  ≈ 6", size=15, bullet=True)
    _add_para(tf, " ", size=10, bullet=False)
    _add_para(tf, "Per-(site, hour) prior buckets: 28 sites × 24 hours", size=15, bullet=True)
    _add_para(tf, "Computed from labels CSV → fed to post-processing", size=14, bullet=False, indent=1)
    _add_footer(s)

    # ---- Slide 7: Proposed Solution Architecture ----
    s = _new_slide(prs)
    _add_title(s, "4. Proposed Solution — Architecture")
    tf = _add_textbox(s, 0.5, 1.4, 12.5, 5.6)
    _add_para(tf, "Pipeline (per soundscape):", size=18, bold=True, bullet=False)
    _add_para(tf, "1. Split 60 s into 12 × 5-s windows", size=16, bullet=True)
    _add_para(tf, "2. Run three independent branches on the same windows:", size=16, bullet=True)
    _add_para(tf, "(a) Perch v2 ONNX → logits + embeddings → mapped to 234 classes", size=15, bullet=False, indent=1)
    _add_para(tf, "(b) Distilled SED 5-fold ONNX → clip + frame_max sigmoid", size=15, bullet=False, indent=1)
    _add_para(tf, "(c) EfficientNetV2-S trained on competition data → BCE logits", size=15, bullet=False, indent=1)
    _add_para(tf, "3. Convert each branch's predictions to per-class column rank-percentiles", size=16, bullet=True)
    _add_para(
        tf, "4. Weighted average:  blend = 0.50·rank_perch + 0.30·rank_sed + 0.20·rank_eff",
        size=16, bullet=True, color=ORANGE,
    )
    _add_para(tf, "5. Post-processing stack (in order):", size=16, bullet=True)
    _add_para(tf, "(a) (site, hour) prior mix at weight 0.05", size=15, bullet=False, indent=1)
    _add_para(tf, "(b) rank-aware scaling: multiply by file_max^0.5", size=15, bullet=False, indent=1)
    _add_para(tf, "(c) delta-shift smoothing: new[t] = 0.85·old[t] + 0.075·(old[t-1] + old[t+1])", size=15, bullet=False, indent=1)
    _add_para(tf, "6. Write 12 × 234-column rows per soundscape to submission.csv", size=16, bullet=True)
    _add_footer(s)

    # ---- Slide 8: Feature Engineering Detail ----
    s = _new_slide(prs)
    _add_title(s, "4.1 Feature Engineering & Class Mapping")
    _add_bullets(
        s,
        [
            "Class mapping (Perch's ~10,000 vocab → 234 target classes):",
            ("Direct match by `scientific_name` join → 203 classes (87%)", 1),
            ("Genus-level proxy (max-pool same-genus Perch logits) → 3 classes", 1),
            ("Unmapped (insect sonotypes, no scientific name) → 28 classes", 1),
            "→ Unmapped classes get Perch's noise floor; EffNet covers them with learned signal",
            "Audio augmentation for EffNet training:",
            ("mixup α=0.4 (p=0.5), random gain ±30%, polarity flip, Gaussian noise σ=0.005", 1),
            ("SpecAugment frequency/time masking (p=0.5)", 1),
            "Cross-platform checkpoint loading: shim WindowsPath → PurePosixPath before torch.load",
            "P100 GPU compatibility: install torch 2.5.1+cu118 in-kernel via nvidia-smi detection",
        ],
        size=17,
    )
    _add_footer(s)

    # ---- Slide 9: Model Justification ----
    s = _new_slide(prs)
    _add_title(s, "4.2 Model Justification — Why Three Branches?")
    _add_table(
        s,
        ["Branch", "Strength", "Weakness", "Role in blend"],
        [
            ("Perch v2 (foundation)", "203/234 classes have direct logits", "Useless on 28 sonotype classes", "Primary signal (weight 0.50)"),
            ("Distilled SED (5 folds)", "Spectrogram-local event evidence", "No transfer learning from foundation", "Complementary (weight 0.30)"),
            ("EffNetV2-S (trained)", "Covers all 234 classes weakly", "Less accurate per-class than Perch on overlap", "Coverage filler (weight 0.20)"),
        ],
        height=2.5,
    )
    _add_bullets(
        s,
        [
            "Rank averaging is asymmetric — each branch contributes where it is informative; others fill in",
            "Heterogeneous input modalities (raw waveform → Perch; mel-spec → SED & EffNet) reduce correlated errors",
            "All inputs are publicly addable per Kaggle code-comp rules",
        ],
        top=4.5, height=2.5, size=16,
    )
    _add_footer(s)

    # ---- Slide 10: LB Progression ----
    s = _new_slide(prs)
    _add_title(s, "5.1 Leaderboard Progression — What Worked")
    _add_table(
        s,
        ["Submission", "Public LB", "Δ from prior"],
        [
            ("EffNetV2-S fold0 only (10 epochs)", "0.856", "—"),
            ("Perch v2 ONNX only (203 + 3 + 28 unmapped)", "0.835", "−0.021"),
            ("3-way rank-average blend", "**0.927**", "**+0.092**"),
            ("+ (site, hour) prior weight 0.05", "0.928", "+0.001"),
            ("+ rank-aware scaling power 0.5", "0.931", "+0.003"),
            ("+ delta-shift smoothing α=0.15  ←  best", "**0.932**", "**+0.001**"),
        ],
        height=4.0,
    )
    _add_bullets(
        s,
        [
            "Single biggest gain (+0.092) came from blending — no individual model is best on its own",
            "Post-processing stack added a further +0.005 on top of the blend baseline",
        ],
        top=6.0, height=1.0, size=15,
    )
    _add_footer(s)

    # ---- Slide 11: Ablations ----
    s = _new_slide(prs)
    _add_title(s, "5.2 Ablations & Sensitivity Analysis")
    _add_table(
        s,
        ["Component change", "LB", "Δ", "Verdict"],
        [
            ("Phase 11 baseline (blend + prior + rank-aware + delta-shift)", "0.932", "—", "kept"),
            ("Tune rank-aware power 0.5 → 0.7", "0.932", "0.000", "flat — reverted"),
            ("Tune (site, hour) prior weight 0.05 → 0.10", "0.932", "0.000", "flat — reverted"),
            ("Replace delta-shift with gaussian σ=0.65", "0.924", "−0.008", "rejected"),
            ("Replace delta-shift with gaussian σ=0.30", "0.928", "−0.004", "rejected"),
            ("Per-class branch weights (Perch=0 on unmapped)", "0.926", "−0.002", "rejected (Section 6)"),
        ],
        height=4.0,
    )
    _add_bullets(
        s,
        [
            "Score plateaus around 0.932 → wide ridge in hyperparameter space",
            "Delta-shift's bounded 3-point kernel beats gaussian with comparable effective radius — sharpness is preserved",
        ],
        top=6.0, height=1.0, size=15,
    )
    _add_footer(s)

    # ---- Slide 12: Failed Experiment 1 — ProtoSSM ----
    s = _new_slide(prs)
    _add_title(s, "6.1 Failed Experiment — ProtoSSM Sequence Head", color=ORANGE)
    _add_bullets(
        s,
        [
            "Hypothesis: bi-LSTM over Perch embeddings would capture inter-window call persistence",
            "What we built: 2-layer bi-LSTM (hidden=256, dropout=0.3), BCE loss, 40 epochs, 66 labeled files",
            "What happened — three failures:",
            ("(P4) In-kernel training exceeded 90-min CPU runtime — submission marked COMPLETE with no score", 1),
            ("(P5) Pre-trained ProtoSSM v1 → blend regressed to 0.908 (−0.019). Loss saturated at 0.052 in 5 of 40 epochs (overfit)", 1),
            ("(P6) Regularized v2 (hidden=128, dropout=0.5, val-split, early stop) → 0.909 (still −0.018 vs blend)", 1),
            "Why it failed:",
            ("(1) 47 effective training files — much smaller than the public 0.943 solution's pseudo-labeled corpus", 1),
            ("(2) Trained with vanilla BCE — the public solution adds knowledge distillation from Perch logits", 1),
            ("(3) Inserted at weight 0.20 without LB-blind validation of its actual contribution", 1),
            "Lesson: a sequence-model branch must (a) be trained with KD against a stronger branch, OR (b) prove validation gain before insertion",
        ],
        size=15,
    )
    _add_footer(s)

    # ---- Slide 13: Failed Experiment 2 — Per-class weights ----
    s = _new_slide(prs)
    _add_title(s, "6.2 Failed Experiment — Per-Class Branch Weights", color=ORANGE)
    _add_bullets(
        s,
        [
            "Hypothesis: Perch outputs constant 0.5 on the 28 unmapped insect classes (pure noise) — drop its weight to 0 for those classes only",
            "Implementation: three weight-vectors keyed by class bucket",
            ("weights        = (0.5, 0.3, 0.2)   # 203 mapped classes", 1),
            ("unmapped_weights = (0.0, 0.5, 0.5)  # 28 sonotype classes", 1),
            ("proxy_weights  = (0.30, 0.40, 0.30) # 3 genus-proxy classes", 1),
            "What happened: LB regressed 0.928 → 0.926 (−0.002)",
            "Why it failed (counterintuitive):",
            ("Perch's constant 0.5 across all rows ties at the median rank — it's not noise in rank space, it's a stabilizer", 1),
            ("Per-class re-weighting changed the rank denominator across classes, disrupting macro-AUC's per-class normalization", 1),
            ("Removing a tied baseline caused subtle relative-rank shifts that hurt rather than helped", 1),
            "Lesson: rank averaging is robust to one branch being uninformative *only when ties are well-distributed*. Per-class weight surgery breaks normalization.",
        ],
        size=15,
    )
    _add_footer(s)

    # ---- Slide 14: Lessons & Conclusion ----
    s = _new_slide(prs)
    _add_title(s, "7. Conclusion & Lessons Learned")
    _add_bullets(
        s,
        [
            "Final result: 0.932 public LB — a 0.097 jump over best single model (0.835), 0.011 below the 0.943 reference",
            "The dominant mechanism for our score is the simple 3-way blend — sophistication of any single branch matters less than coverage diversity",
            "Lessons (in priority order):",
            ("1. Validation infrastructure is the bottleneck. Without 5-fold OOF stacking, every hyperparameter is a blind LB shot. The 5/day quota cap is brutal.", 1),
            ("2. Inserting a new branch into a strong blend is risky. Even a well-trained branch can hurt if it lacks KD or coverage diversity.", 1),
            ("3. Surgical class-level corrections can disturb metric normalization. Rank averaging tolerates noise better than tampering.", 1),
            ("4. Foundation models do most of the work. 87% of the score comes from a 1-line scientific_name join.", 1),
            "Future work to close the 0.011 gap:",
            ("Build OOF infrastructure → enables per-class temperature, per-taxon thresholds, blend-weight tuning", 1),
            ("Proper ProtoSSM with KD from Perch + mixup + focal + SWA → +0.005 to +0.010 expected", 1),
            ("Iterative pseudo-labeling on unlabeled train_soundscapes → +0.005 to +0.010 expected", 1),
        ],
        size=16,
    )
    _add_footer(s)

    # ---- Slide 15: References ----
    s = _new_slide(prs)
    _add_title(s, "References & Acknowledgements")
    _add_bullets(
        s,
        [
            "Kaggle, 'BirdCLEF+ 2026 — Acoustic Species Identification', 2026. kaggle.com/competitions/birdclef-2026",
            "Google Research + Cornell Lab of Ornithology — Perch v2 (foundation model)",
            "Tucker Arrants — public distilled SED 5-fold ONNX ensemble",
            "Rishikesh Jani — Perch ONNX export for BirdCLEF+ 2026",
            "Mattia Angeli — public 0.943 ProtoSSM+SED notebook (architecture reference)",
            "Vyanktesh Dwivedi — public ONNX-Perch sequence-modeling base notebook",
            "Jaejohn — perch-meta cached embeddings (used for ProtoSSM training)",
            "Ross Wightman — `timm` PyTorch image models library",
            " ",
            "Code: github.com/RamananVr/KaggleBirdClef2026",
            "Tag for the 0.927 architecture baseline:  v0.927-baseline",
        ],
        size=14,
    )
    _add_footer(s)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}")


if __name__ == "__main__":
    build()
